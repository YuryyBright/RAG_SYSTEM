
from pathlib import Path
from .base_reader import BaseReader

class PptxReader(BaseReader):
    """
    Reader for Microsoft PowerPoint (.pptx) files.

    Reads all text from the slides in a .pptx file.

    Parameters
    ----------
    path : Path
        Path to the .pptx file.

    Returns
    -------
    str
        Text content from the .pptx file.
    """

    def read(self, path: Path) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)
        except Exception as e:
            return f"Error reading PPTX file: {str(e)}"

