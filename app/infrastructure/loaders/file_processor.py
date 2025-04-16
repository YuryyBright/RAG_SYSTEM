# app/infrastructure/loaders/file_processor.py
import os
import uuid
import json
import logging
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime
from pathlib import Path
from dataclasses import asdict
from langdetect import detect
from langdetect.lang_detect_exception import LangDetectException
from sqlalchemy.ext.asyncio import AsyncSession

from core.entities.document import Document
from core.entities.processed_file import ProcessedFile
from infrastructure.database.repository.file_repository import FileRepository
from utils.logger_util import get_logger

logger = get_logger(__name__)


class FileProcessor:
    """
    File processor for reading various file types, detecting language,
    and preparing text for vector database storage in a RAG system.

    Provides detailed reporting on processed files, including warnings for
    unreadable files and language detection failures.
    """

    def __init__(self, file_repository: FileRepository):

        self.supported_extensions = {
            "txt", "md", "markdown", "html", "htm", "pdf", "doc", "docx",
            "csv", "xls", "xlsx", "json", "xml", "ppt", "pptx"
        }
        self.successful_files = []
        self.unreadable_files = []
        self.language_detection_failures = []
        self.files_with_warnings = []

        # Map extensions to their handler methods
        self.extension_handlers = {
            ".txt": self._read_text_file,
            ".md": self._read_text_file,
            ".json": self._read_json_file,
            ".html": self._read_html_file,
            ".csv": self._read_csv_file,
            ".pdf": self._read_pdf_file,
            ".docx": self._read_docx_file,
            ".xlsx": self._read_xlsx_file,
            ".pptx": self._read_pptx_file,
            ".rtf": self._read_rtf_file,
            ".xml": self._read_xml_file,
            ".yaml": self._read_yaml_file,
            ".yml": self._read_yaml_file,
        }
        self.file_repository = file_repository
    async def process_directory(
            self,
            directory_path: str,
            recursive: bool = True,
            metadata: Optional[Dict[str, Any]] = None
    ) -> Tuple[List[Document], Dict[str, Any]]:
        """
        Process all files in a directory for the RAG system.

        Parameters
        ----------
        directory_path : str
            Path to the directory containing files to process.
        recursive : bool, optional
            Whether to process subdirectories (default is True).
        metadata : Dict[str, Any], optional
            Additional metadata to attach to all documents.

        Returns
        -------
        Tuple[List[Document], Dict[str, Any]]
            - List of processed documents ready for vector database
            - Processing report with details on successful, unreadable, and warning files
        """
        path = Path(directory_path)
        if not path.exists() or not path.is_dir():
            raise ValueError(f"Directory {directory_path} does not exist or is not a directory")

        # Clear previous processing results
        self.successful_files.clear()
        self.unreadable_files.clear()
        self.language_detection_failures.clear()
        self.files_with_warnings.clear()

        # Get all files or files in all subdirectories if recursive
        files = path.rglob("*") if recursive else path.glob("*")

        documents = []
        for file_path in files:
            if not file_path.is_file():
                continue

            suffix = file_path.suffix.lower()

            # Skip unsupported file types
            if suffix not in self.extension_handlers:
                logger.info(f"Skipping unsupported file type: {file_path}")
                continue

            try:
                processed_file = await self.process_file(file_path, metadata)

                # Handle the processed file based on its status
                if processed_file.is_readable:
                    doc = await self._convert_to_document(processed_file, file_path)
                    documents.append(doc)

                    has_warnings = False

                    # Check for language detection failures
                    if not processed_file.is_language_detected:
                        self.language_detection_failures.append(processed_file)
                        has_warnings = True

                    # Check for other warnings in metadata
                    if processed_file.metadata.get("has_warnings", False):
                        has_warnings = True

                    if has_warnings:
                        self.files_with_warnings.append(processed_file)
                    else:
                        self.successful_files.append(processed_file)
                else:
                    self.unreadable_files.append(processed_file)
                logger.info(f"Finished processing file {file_path}")
            except Exception as e:
                logger.error(f"Failed to process {file_path}: {e}")
                # Create a processed file record for the failure
                failed_file = ProcessedFile(
                    id=str(uuid.uuid4()),
                    filename=file_path.name,
                    content="",
                    is_readable=False,
                    metadata={
                        "source": str(file_path),
                        "error": str(e),
                        "file_size": file_path.stat().st_size if file_path.exists() else 0
                    }
                )
                self.unreadable_files.append(failed_file)

        # Generate the report
        report = self._generate_enhanced_report()
        return documents, report

    async def process_file(self, file_path: Path, metadata: Optional[Dict[str, Any]] = None) -> ProcessedFile:
        """
        Process a single file: read content and detect language.

        Parameters
        ----------
        file_path : Path
            Path to the file to process.
        metadata : Dict[str, Any], optional
            Additional metadata to attach to the processed file.

        Returns
        -------
        ProcessedFile
            Processed file entity with content and metadata.
        """
        # Create base metadata
        meta = {
            "source": str(file_path),
            "filename": file_path.name,
            "extension": file_path.suffix,
            "file_size": file_path.stat().st_size,
            "created_at": datetime.fromtimestamp(file_path.stat().st_ctime).isoformat(),
            "modified_at": datetime.fromtimestamp(file_path.stat().st_mtime).isoformat(),
            "processed_at": datetime.now().isoformat(),
            "has_warnings": False
        }

        # Add user-provided metadata
        if metadata:
            meta.update(metadata)

        content = ""
        is_readable = True

        try:
            # Get the appropriate handler for this file extension
            handler = self.extension_handlers.get(file_path.suffix.lower())
            if handler:
                content = handler(file_path)
            else:
                is_readable = False
                meta["error"] = f"No handler for file extension: {file_path.suffix}"
                meta["has_warnings"] = True

        except Exception as e:
            is_readable = False
            meta["error"] = str(e)
            meta["has_warnings"] = True
            logger.warning(f"Could not read file {file_path}: {e}")

        # Detect language if content is readable
        language = None
        is_language_detected = False

        if is_readable and content.strip():
            try:
                # Only attempt language detection if content is substantial
                if len(content.strip()) > 20:  # Increased minimum content length for better detection
                    language = detect(content[:5000])  # Use only first part for efficiency
                    is_language_detected = True
                else:
                    meta["language_error"] = "Content too short for reliable language detection"
                    meta["has_warnings"] = True
            except LangDetectException as e:
                meta["language_error"] = f"Language detection failed: {str(e)}"
                meta["has_warnings"] = True
            except Exception as e:
                meta["language_error"] = f"Error during language detection: {str(e)}"
                meta["has_warnings"] = True

        # Create processed file entity
        return ProcessedFile(
            id=str(uuid.uuid4()),
            filename=file_path.name,
            content=content,
            language=language,
            is_readable=is_readable,
            is_language_detected=is_language_detected,
            metadata=meta
        )

    async def _convert_to_document(self, pf: ProcessedFile, file_path) -> Document:
        """
        Convert a ProcessedFile to a Document for vector database storage.

        Parameters
        ----------
        pf : ProcessedFile
            The processed file to convert.

        Returns
        -------
        Document
            Document entity ready for vector database.
        """
        # Copy metadata to avoid modifying the original
        meta = pf.metadata.copy() if pf.metadata else {}

        # Add language info to metadata
        if pf.language:
            meta["language"] = pf.language

        # Add warning flags
        if not pf.is_language_detected:
            meta["language_detection_failed"] = True

        # Fetch file from DB using path
        try:
            file_record = await self.file_repository.get_file_by_path((str(file_path)))
            if file_record:
                return Document(id=pf.id, content=pf.content, metadata=meta, file_id=file_record.id, owner_id=file_record.owner_id)
            raise "File not found"
        except Exception as e:
            logger.warning(f"Failed to fetch file from DB for path {meta.get('source')}: {e}")



    def _generate_enhanced_report(self) -> Dict[str, Any]:
        """
        Generate a detailed report of file processing results.

        Returns
        -------
        Dict[str, Any]
            Processing report with summary, details, and recommendations.
        """
        # Extract data for the report
        successful = [asdict(f) for f in self.successful_files]
        unreadable = [asdict(f) for f in self.unreadable_files]
        language_failures = [asdict(f) for f in self.language_detection_failures]
        warnings = [asdict(f) for f in self.files_with_warnings]

        # Files that might need manual review
        files_to_review = []
        for file in self.files_with_warnings:
            files_to_review.append(file.filename)

        # Files that might need to be removed
        files_to_consider_removing = []
        for file in self.unreadable_files:
            files_to_consider_removing.append(file.filename)

        # For empty files with no errors, suggest removal
        for file in self.successful_files:
            if not file.content.strip() and file.is_readable:
                files_to_consider_removing.append(file.filename)

        return {
            "summary": {
                "total_files": len(successful) + len(unreadable) + len(warnings),
                "successful": len(successful),
                "unreadable": len(unreadable),
                "language_detection_failures": len(language_failures),
                "files_with_warnings": len(warnings)
            },
            "details": {
                "successful_files": [{"filename": f["filename"], "language": f["language"]} for f in successful],
                "unreadable_files": [{"filename": f["filename"], "error": f["metadata"].get("error")} for f in
                                     unreadable],
                "language_detection_failures": [
                    {"filename": f["filename"], "error": f["metadata"].get("language_error")} for f in
                    language_failures],
                "files_with_warnings": [{"filename": f["filename"], "warnings": self._extract_warnings(f["metadata"])}
                                        for f in warnings]
            },
            "recommendations": {
                "files_to_review": files_to_review,
                "files_to_consider_removing": files_to_consider_removing
            }
        }

    def _extract_warnings(self, metadata: Dict[str, Any]) -> List[str]:
        """
        Extract warning messages from file metadata.

        Parameters
        ----------
        metadata : Dict[str, Any]
            File metadata containing warning information.

        Returns
        -------
        List[str]
            List of warning messages.
        """
        warnings = []
        if metadata.get("language_error"):
            warnings.append(f"Language detection issue: {metadata['language_error']}")
        if metadata.get("error"):
            warnings.append(f"Processing error: {metadata['error']}")

        return warnings

    # === File readers ===
    def _read_text_file(self, path: Path) -> str:
        """Read content from text files."""
        return path.read_text(encoding="utf-8", errors="replace")

    def _read_json_file(self, path: Path) -> str:
        """Read and extract text from JSON files."""
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            try:
                data = json.load(f)
                if isinstance(data, dict):
                    # Try to extract content field or stringify the dict
                    return data.get("content", json.dumps(data, ensure_ascii=False))
                elif isinstance(data, list):
                    # Stringify each item in the list
                    return "\n".join(json.dumps(item, ensure_ascii=False) for item in data)
                return str(data)
            except json.JSONDecodeError as e:
                # Handle invalid JSON
                return f"Invalid JSON file: {str(e)}"

    def _read_html_file(self, path: Path) -> str:
        """Extract text content from HTML files."""
        try:
            from bs4 import BeautifulSoup
            content = path.read_text(encoding="utf-8", errors="replace")
            soup = BeautifulSoup(content, "html.parser")
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.extract()
            # Get text
            return soup.get_text(separator=" ", strip=True)
        except ImportError:
            # Fallback to simple HTML parser if BeautifulSoup is not available
            from html.parser import HTMLParser
            class SimpleHTMLParser(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.result = []

                def handle_data(self, d):
                    self.result.append(d)

            parser = SimpleHTMLParser()
            parser.feed(path.read_text(encoding="utf-8", errors="replace"))
            return " ".join(parser.result)

    def _read_csv_file(self, path: Path) -> str:
        """Read and format text from CSV files."""
        import csv
        try:
            with open(path, newline='', encoding='utf-8', errors="replace") as f:
                reader = csv.reader(f)
                return "\n".join([", ".join(row) for row in reader])
        except Exception as e:
            return f"Error reading CSV: {str(e)}"

    def _read_pdf_file(self, path: Path) -> str:
        """Extract text from PDF files."""
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(path)
            text = ""
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            return text
        except Exception as e:
            return f"Error reading PDF: {str(e)}"

    def _read_docx_file(self, path: Path) -> str:
        """Extract text from DOCX files."""
        try:
            from docx import Document
            doc = Document(path)
            return "\n".join(paragraph.text for paragraph in doc.paragraphs)
        except Exception as e:
            return f"Error reading DOCX: {str(e)}"

    def _read_xlsx_file(self, path: Path) -> str:
        """Extract text from XLSX files."""
        try:
            from openpyxl import load_workbook
            wb = load_workbook(path, data_only=True)
            result = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    result.append(", ".join(str(cell.value) if cell.value is not None else "" for cell in row))
            return "\n".join(result)
        except Exception as e:
            return f"Error reading XLSX: {str(e)}"

    def _read_pptx_file(self, path: Path) -> str:
        """Extract text from PPTX files."""
        try:
            from pptx import Presentation
            prs = Presentation(path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            return f"Error reading PPTX: {str(e)}"

    def _read_rtf_file(self, path: Path) -> str:
        """Extract text from RTF files."""
        try:
            import striprtf.striprtf as striprtf
            with open(path, 'rb') as f:
                rtf_content = f.read().decode('utf-8', errors="replace")
                return striprtf.rtf_to_text(rtf_content)
        except Exception as e:
            return f"Error reading RTF: {str(e)}"

    def _read_xml_file(self, path: Path) -> str:
        """Extract text from XML files."""
        try:
            import xml.etree.ElementTree as ET
            tree = ET.parse(path)
            root = tree.getroot()

            # Function to extract text from element and its children
            def extract_text(element):
                text = element.text or ""
                for child in element:
                    text += extract_text(child)
                if element.tail:
                    text += element.tail
                return text

            return extract_text(root)
        except Exception as e:
            return f"Error reading XML: {str(e)}"

    def _read_yaml_file(self, path: Path) -> str:
        """Extract text from YAML files."""
        try:
            import yaml
            with open(path, 'r', encoding='utf-8', errors="replace") as f:
                data = yaml.safe_load(f)
                return yaml.dump(data, default_flow_style=False)
        except Exception as e:
            return f"Error reading YAML: {str(e)}"